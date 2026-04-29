"""
PowerPoint Presentation Generator
Creates professional slides with analysis, findings, and resources
"""

import logging
from typing import List, Dict
from datetime import datetime

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("[PPT] python-pptx not installed. Install with: pip install python-pptx")


class PowerPointGenerator:
    """
    Generate professional PowerPoint presentations with:
    - Title slide
    - Executive summary
    - Key findings
    - Paper summaries (6 slides)
    - Trend analysis
    - Call-to-action
    """

    def __init__(self, output_path: str = "results/report.pptx"):
        """Initialize PowerPoint generator"""
        self.output_path = output_path
        self.has_pptx = HAS_PPTX
        self.colors = {
            'primary': RGBColor(102, 126, 234),      # #667eea
            'secondary': RGBColor(118, 75, 162),     # #764ba2
            'heading': RGBColor(26, 32, 44),         # #1a202c
            'text': RGBColor(45, 55, 72),            # #2d3748
            'accent': RGBColor(237, 137, 54),        # #ed8936
            'white': RGBColor(255, 255, 255),
        }

        if not self.has_pptx:
            logger.error("[PPT] python-pptx required. Install: pip install python-pptx")

    def _add_hyperlink_run(self, paragraph, text: str, url: str,
                           font_size=Pt(11), color=None, bold=False, underline=True):
        """
        Add a clickable hyperlink run to a paragraph.
        Uses python-pptx's proper hyperlink API for real clickable links in PowerPoint.
        """
        if color is None:
            color = RGBColor(0x33, 0x6B, 0xB4)  # Professional blue

        run = paragraph.add_run()
        run.text = text
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.underline = underline

        if url:
            # Normalize URL
            if not url.startswith(('http://', 'https://')):
                url = 'https://' + url
            run.hyperlink.address = url

        return run

    def generate(
        self,
        insights: List[Dict],
        all_sources: List[Dict] = None,
        report_bundle: Dict = None,
    ) -> bool:
        """
        Generate comprehensive PowerPoint presentation
        Returns: True if successful
        """
        if not self.has_pptx:
            logger.error("[PPT] Cannot generate - python-pptx not installed")
            return False

        if not insights:
            logger.warning("[PPT] No insights to generate presentation")
            return False

        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Build slides
            self._add_title_slide(prs, insights)
            self._add_executive_summary_slide(prs, insights)
            self._add_key_findings_slide(prs, insights)

            # ── HARNESS PRINCIPLE: Source-Grouped Selection ──────────────
            # Group by source, pick #1 from each, then fill remaining slots.
            harness_insights = self._harness_select(insights)

            # Add paper slides (1 top per source)
            for idx, paper in enumerate(harness_insights, 1):
                self._add_paper_slide(prs, paper, idx)

            self._add_trends_slide(prs, insights)

            # Add selected sources slide (only harness-selected, with clickable links)
            self._add_selected_sources_slide(prs, harness_insights)

            self._add_cta_slide(prs)

            # Save
            prs.save(self.output_path)
            logger.info(f"[PPT] Presentation generated: {self.output_path}")
            return True

        except Exception as e:
            logger.error(f"[PPT] Generation failed: {e}")
            return False

    # ── Harness Principle ─────────────────────────────────────────────────

    def _harness_select(
        self, insights: List[Dict], top_per_source: int = 1, max_slides: int = 8
    ) -> List[Dict]:
        """
        Source-grouped selection: pick the #1 item from each distinct source,
        then fill remaining slots with the next-best globally.

        This ensures balanced coverage across Apple, DeepMind, arXiv, Meta, etc.
        instead of all slides coming from a single dominant source.
        """
        from collections import defaultdict

        # Group by normalised source name
        by_source = defaultdict(list)
        for item in insights:
            src = (
                item.get('source', '')
                or item.get('source_platform', '')
                or 'Unknown'
            ).strip().lower()
            by_source[src].append(item)

        # Sort each group by score (descending)
        for src in by_source:
            by_source[src].sort(
                key=lambda x: float(x.get('relevance_score', 0)), reverse=True
            )

        selected = []
        used_titles = set()

        # Pass 1: top N from each source
        for src, items in by_source.items():
            for item in items[:top_per_source]:
                title = item.get('title', '')
                if title not in used_titles:
                    selected.append(item)
                    used_titles.add(title)

        # Pass 2: fill remaining slots from globally sorted leftovers
        if len(selected) < max_slides:
            all_sorted = sorted(
                insights,
                key=lambda x: float(x.get('relevance_score', 0)),
                reverse=True,
            )
            for item in all_sorted:
                if len(selected) >= max_slides:
                    break
                title = item.get('title', '')
                if title not in used_titles:
                    selected.append(item)
                    used_titles.add(title)

        # Final sort so slides appear in score order
        selected.sort(
            key=lambda x: float(x.get('relevance_score', 0)), reverse=True
        )

        logger.info(
            f"[PPT] Harness selected {len(selected)} papers from "
            f"{len(by_source)} sources (max {max_slides})"
        )
        return selected

    def _add_title_slide(self, prs: Presentation, insights: List[Dict]):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "On-Device AI Intelligence"
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['white']

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Research Intelligence Report • {datetime.now().strftime('%B %d, %Y')}"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = self.colors['white']

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{len(insights)} Papers Analyzed • Hybrid RAG + Multi-Model AI"
        footer_frame.paragraphs[0].font.size = Pt(16)
        footer_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    def _add_executive_summary_slide(self, prs: Presentation, insights: List[Dict]):
        """Add executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Executive Summary"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Calculate metrics
        total = len(insights)
        avg_score = sum(i.get('relevance_score', 0) for i in insights) / total if total > 0 else 0

        platforms = {}
        for item in insights:
            platform = item.get('platform', 'Unknown')
            platforms[platform] = platforms.get(platform, 0) + 1

        high_impact = len([i for i in insights if i.get('dram_impact') == 'High'])

        # Content
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        metrics = [
            f"📊 Total Papers: {total}",
            f"⭐ Average Score: {avg_score:.1f}/100",
            f"📱 Mobile Papers: {platforms.get('Mobile', 0)}",
            f"💻 Laptop Papers: {platforms.get('Laptop', 0)}",
            f"🔥 High Impact: {high_impact} papers",
        ]

        for idx, metric in enumerate(metrics):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = metric
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['text']
            p.space_before = Pt(12)

    def _add_key_findings_slide(self, prs: Presentation, insights: List[Dict]):
        """Add key findings slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Key Findings"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Techniques
        techniques = {}
        for item in insights:
            tech = item.get('quantization_method', 'N/A')
            if tech != 'N/A':
                techniques[tech] = techniques.get(tech, 0) + 1

        top_techniques = sorted(techniques.items(), key=lambda x: x[1], reverse=True)[:3]

        # Content
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        findings = ["🎯 Top Optimization Techniques:"]
        for tech, count in top_techniques:
            findings.append(f"   • {tech}: {count} papers")

        findings.append("")

        # Data-driven insights instead of hardcoded generic text
        dram_dist = {}
        platform_dist = {}
        for item in insights:
            d = item.get('dram_impact', 'Unknown')
            dram_dist[d] = dram_dist.get(d, 0) + 1
            p = item.get('platform', 'Unknown')
            platform_dist[p] = platform_dist.get(p, 0) + 1

        findings.append("📊 DRAM Impact Distribution:")
        for impact, count in sorted(dram_dist.items(), key=lambda x: x[1], reverse=True):
            findings.append(f"   • {impact}: {count} papers")

        findings.append("")
        findings.append("🖥️ Platform Coverage:")
        for plat, count in sorted(platform_dist.items(), key=lambda x: x[1], reverse=True)[:4]:
            findings.append(f"   • {plat}: {count} papers")

        # Top paper highlight
        top = sorted(insights, key=lambda x: x.get('relevance_score', 0), reverse=True)
        if top:
            top_title = top[0].get('title', 'Unknown')[:60]
            top_score = top[0].get('relevance_score', 0)
            findings.append("")
            findings.append(f"⭐ Top Paper: \"{top_title}\" (Score: {top_score}/100)")

        for idx, finding in enumerate(findings):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = finding
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.space_before = Pt(4)

    def _add_paper_slide(self, prs: Presentation, paper: Dict, rank: int):
        """Add structured research brief slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        paper_title = paper.get('title', 'Unknown')
        title.text = f"#{rank}: {paper_title[:55]}"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # ── Metadata ribbon ──────────────────────────────────────────────
        score = paper.get('relevance_score', 0)
        platform = paper.get('platform', 'Unknown')
        model_type = paper.get('model_type', 'Unknown')
        dram_impact = paper.get('dram_impact', 'Unknown')
        source = paper.get('source', 'Unknown')
        paper_url = paper.get('link', '') or paper.get('url', '')

        ribbon_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.4))
        ribbon_frame = ribbon_box.text_frame
        p = ribbon_frame.paragraphs[0]
        # Add metadata as plain text run
        meta_run = p.add_run()
        meta_run.text = (
            f"Score: {score}/100  •  {platform}  •  {model_type}  •  "
            f"DRAM: {dram_impact}  •  "
        )
        meta_run.font.size = Pt(11)
        meta_run.font.color.rgb = RGBColor(0x88, 0x88, 0xAA)
        meta_run.font.bold = True

        # Add source as clickable hyperlink
        if paper_url:
            self._add_hyperlink_run(
                p, f"📎 {source}", paper_url,
                font_size=Pt(11), color=RGBColor(0x33, 0x6B, 0xB4), bold=True
            )
        else:
            src_run = p.add_run()
            src_run.text = f"Source: {source}"
            src_run.font.size = Pt(11)
            src_run.font.color.rgb = RGBColor(0x88, 0x88, 0xAA)
            src_run.font.bold = True

        # ── Main content area ────────────────────────────────────────────
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(9), Inches(5.2))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Try new structured fields first, fall back to legacy
        exec_summary = paper.get('executive_summary', '')
        research_brief = paper.get('research_brief', {})
        key_metrics = paper.get('key_metrics', [])
        why_matters = paper.get('why_it_matters', '')
        impl_guide = paper.get('implementation_guidance', '')

        # Legacy fallback
        if not exec_summary:
            memory_insight = str(paper.get('memory_insight', 'N/A'))
            takeaway = str(paper.get('engineering_takeaway', 'N/A'))
            exec_summary = f"{memory_insight} {takeaway}"

        content_blocks = []

        # Executive Summary
        content_blocks.append(("📋 Executive Summary", True, Pt(13), self.colors['accent']))
        content_blocks.append((exec_summary[:400], False, Pt(12), self.colors['text']))
        content_blocks.append(("", False, Pt(6), self.colors['text']))

        # Research Brief (Problem → Method → Result)
        if isinstance(research_brief, dict) and research_brief:
            content_blocks.append(("🔬 Research Brief", True, Pt(13), self.colors['accent']))
            problem = research_brief.get('problem', '')
            method = research_brief.get('method', '')
            result = research_brief.get('result', '')
            if problem:
                content_blocks.append((f"Problem: {problem[:150]}", False, Pt(11), self.colors['text']))
            if method:
                content_blocks.append((f"Method: {method[:150]}", False, Pt(11), self.colors['text']))
            if result:
                content_blocks.append((f"Result: {result[:150]}", False, Pt(11), self.colors['text']))
            content_blocks.append(("", False, Pt(6), self.colors['text']))

        # Key Metrics
        if key_metrics and isinstance(key_metrics, list):
            content_blocks.append(("📊 Key Metrics", True, Pt(13), self.colors['accent']))
            metrics_text = "  •  ".join(str(m)[:60] for m in key_metrics[:5])
            content_blocks.append((metrics_text, False, Pt(11), self.colors['text']))
            content_blocks.append(("", False, Pt(6), self.colors['text']))

        # Why It Matters
        if why_matters:
            content_blocks.append(("💡 Why It Matters", True, Pt(13), self.colors['accent']))
            content_blocks.append((str(why_matters)[:250], False, Pt(11), self.colors['text']))

        # Implementation Guidance
        if impl_guide:
            content_blocks.append(("⚙️ Implementation", True, Pt(12), self.colors['accent']))
            content_blocks.append((str(impl_guide)[:200], False, Pt(11), self.colors['text']))

        # Render all blocks
        for idx, (text, bold, size, color) in enumerate(content_blocks):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = text
            p.font.size = size
            p.font.color.rgb = color
            p.font.bold = bold
            p.space_before = Pt(2)

    def _add_trends_slide(self, prs: Presentation, insights: List[Dict]):
        """Add trends slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Research Trends"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Sources
        sources = {}
        for item in insights:
            source = item.get('source', 'Unknown')
            sources[source] = sources.get(source, 0) + 1

        top_sources = sorted(sources.items(), key=lambda x: x[1], reverse=True)[:3]

        # Content
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        trends = ["📈 Most Active Sources:"]
        for source, count in top_sources:
            trends.append(f"   • {source}: {count} papers")

        trends.append("")
        trends.append("🔮 Emerging Patterns:")
        trends.append("   • Increasing focus on edge device optimization")
        trends.append("   • Memory efficiency is the top priority")
        trends.append("   • Cross-platform compatibility studies growing")

        for idx, trend in enumerate(trends):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = trend
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']
            p.space_before = Pt(6)

    def _add_selected_sources_slide(self, prs: Presentation, selected_papers: List[Dict]):
        """
        Add a clean 'Selected Sources' slide showing ONLY the harness-selected
        papers with clickable hyperlinks. No bloat — just the curated picks.
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"Selected Sources ({len(selected_papers)} papers)"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.4))
        sub_frame = sub_box.text_frame
        sub_frame.text = "Click any title to open the original paper"
        sub_frame.paragraphs[0].font.size = Pt(12)
        sub_frame.paragraphs[0].font.color.rgb = RGBColor(0x88, 0x88, 0xAA)
        sub_frame.paragraphs[0].font.italic = True

        # Sources list
        sources_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(5.2))
        sources_frame = sources_box.text_frame
        sources_frame.word_wrap = True

        for idx, paper in enumerate(selected_papers):
            if idx > 0:
                sources_frame.add_paragraph()

            p = sources_frame.paragraphs[idx]
            paper_title = paper.get('title', 'Unknown')[:65]
            paper_url = paper.get('link', '') or paper.get('url', '')
            score = paper.get('relevance_score', 0)
            source = paper.get('source', 'Unknown')
            platform = paper.get('platform', 'Unknown')

            # Number prefix
            num_run = p.add_run()
            num_run.text = f"{idx + 1}. "
            num_run.font.size = Pt(12)
            num_run.font.bold = True
            num_run.font.color.rgb = self.colors['accent']

            # Paper title as clickable hyperlink
            if paper_url:
                self._add_hyperlink_run(
                    p, paper_title, paper_url,
                    font_size=Pt(12), bold=False
                )
            else:
                title_run = p.add_run()
                title_run.text = paper_title
                title_run.font.size = Pt(12)
                title_run.font.color.rgb = self.colors['text']

            # Metadata suffix
            meta_run = p.add_run()
            meta_run.text = f"  ({source} • {platform} • Score: {score}/100)"
            meta_run.font.size = Pt(10)
            meta_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

            p.space_before = Pt(6)

        logger.info(f"[PPT] Selected sources slide: {len(selected_papers)} clickable links")

    def _build_fallback_sources(self, insights: List[Dict]) -> List[Dict]:
        """
        Fallback method to extract sources if SourceLinkProcessor unavailable
        """
        sources = []
        seen_urls = set()

        for paper in insights:
            url = paper.get('link') or paper.get('url', '')
            if not url or url in seen_urls:
                continue

            # Normalize URL
            if url and not url.startswith(('http://', 'https://')):
                url = 'https://' + url

            seen_urls.add(url)
            sources.append({
                'title': paper.get('title', 'Unknown Title'),
                'url': url,
                'source_platform': paper.get('source', 'Unknown'),
                'relevance_score': float(paper.get('relevance_score', 0)),
                'summary': paper.get('summary', '')
            })

        # Sort by relevance
        sources.sort(key=lambda x: x['relevance_score'], reverse=True)
        return sources

    def _add_cta_slide(self, prs: Presentation):
        """Add call-to-action slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['secondary']

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "Next Steps"
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['white']

        # Action items
        actions_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3))
        actions_frame = actions_box.text_frame
        actions_frame.word_wrap = True

        actions = [
            "📄 Download the full PDF report",
            "🎤 Listen to the audio podcast",
            "📊 View detailed analysis dashboard",
            "🔗 Access all paper resources",
        ]

        for idx, action in enumerate(actions):
            if idx > 0:
                actions_frame.add_paragraph()
            p = actions_frame.paragraphs[idx]
            p.text = action
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']
            p.space_before = Pt(10)


__all__ = ['PowerPointGenerator']
